"""
Build aibrd product deck as a .pptx file.
Run: python3 scripts/build_ppt.py
Output: aibrd-product-deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x11, 0x17)   # near-black navy
ACCENT     = RGBColor(0x3B, 0x82, 0xF6)   # blue-500
ACCENT2    = RGBColor(0x10, 0xB9, 0x81)   # emerald-500
ACCENT3    = RGBColor(0xF5, 0x9E, 0x0B)   # amber-500
ACCENT4    = RGBColor(0xEF, 0x44, 0x44)   # red-500
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GREY   = RGBColor(0x47, 0x55, 0x69)
CARD_BG    = RGBColor(0x1E, 0x29, 0x3B)   # slate-800

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height,
        fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def label(slide, text, left, top, width, height,
          size=Pt(14), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def multiline(slide, lines, left, top, width, height,
              size=Pt(13), color=LIGHT_GREY, line_spacing=Pt(6),
              bold_first=False):
    """Add a textbox with multiple paragraphs."""
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = (i == 0 and bold_first)
    return txb

def accent_bar(slide, color=ACCENT, height=0.04):
    """Thin colored bar at the top of a slide."""
    bar = box(slide, 0, 0, 13.33, height, fill_color=color)
    return bar

def section_pill(slide, text, left, top, color=ACCENT):
    """Small colored pill label."""
    pill = box(slide, left, top, 1.8, 0.28, fill_color=color)
    pill.line.fill.background()
    label(slide, text, left + 0.05, top + 0.03, 1.7, 0.22,
          size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def divider(slide, top, color=MID_GREY):
    bar = box(slide, 0.5, top, 12.33, 0.01, fill_color=color)

def card(slide, left, top, width, height, title, body_lines,
         title_color=WHITE, body_color=LIGHT_GREY,
         card_color=CARD_BG, accent_color=ACCENT):
    """Rounded-feel card with colored left border."""
    card_shape = box(slide, left, top, width, height, fill_color=card_color)
    accent_shape = box(slide, left, top, 0.05, height, fill_color=accent_color)
    label(slide, title, left + 0.15, top + 0.08, width - 0.2, 0.3,
          size=Pt(11), bold=True, color=title_color)
    body_top = top + 0.38
    for line in body_lines:
        label(slide, line, left + 0.15, body_top, width - 0.25, 0.25,
              size=Pt(9.5), color=body_color)
        body_top += 0.22
    return card_shape


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)

# Background gradient feel — two overlapping boxes
box(sl, 0, 0, 7, 7.5, fill_color=RGBColor(0x0D, 0x11, 0x17))
box(sl, 7, 0, 6.33, 7.5, fill_color=RGBColor(0x13, 0x1C, 0x2E))

# Blue accent bar
box(sl, 0, 0, 13.33, 0.06, fill_color=ACCENT)

# Tag line pill
section_pill(sl, "v0.2.0  ·  MIT License", 0.5, 0.5, color=MID_GREY)

# Title
label(sl, "aibrd", 0.5, 1.1, 9, 1.4,
      size=Pt(72), bold=True, color=WHITE)

# Subtitle
label(sl, "BRD → Living Specification → Code Traceability",
      0.5, 2.55, 9, 0.6, size=Pt(22), color=ACCENT, bold=True)

# Description
label(sl,
      "Converts Business Requirements Documents into a version-controlled\n"
      "living specification inside your repository — keeping the thread\n"
      "between what the business asked for and what shipped, alive permanently.",
      0.5, 3.3, 8, 1.4, size=Pt(14), color=LIGHT_GREY)

# Two mode pills
for i, (txt, col) in enumerate([
    ("VS Code Extension  ·  No API key — uses Copilot", ACCENT),
    ("Python Library  ·  Bring your own LLM key", ACCENT2),
]):
    y = 4.85 + i * 0.55
    box(sl, 0.5, y, 7.5, 0.42, fill_color=col)
    label(sl, txt, 0.65, y + 0.07, 7.3, 0.3,
          size=Pt(12), bold=True, color=WHITE)

# Right side — animated code snippet feel
code_lines = [
    "$ aibrd init ./brd.pdf",
    "",
    "  Detecting modules...",
    "  Extracting: payments...",
    "  Extracting: auth...",
    "",
    "  ✓ 12 flows  ·  24 rules  ·  18 AC",
    "  ✓ .aibrd/ written — commit to git",
]
box(sl, 7.5, 0.8, 5.5, 3.8, fill_color=RGBColor(0x0F, 0x17, 0x2A))
box(sl, 7.5, 0.8, 5.5, 0.28, fill_color=RGBColor(0x1E, 0x29, 0x3B))
for dot_x, dot_c in [(7.65, RGBColor(0xEF,0x44,0x44)),
                      (7.85, RGBColor(0xF5,0x9E,0x0B)),
                      (8.05, RGBColor(0x10,0xB9,0x81))]:
    box(sl, dot_x, 0.9, 0.12, 0.08, fill_color=dot_c)
for i, line in enumerate(code_lines):
    col = ACCENT2 if line.startswith("  ✓") else (ACCENT if line.startswith("$") else LIGHT_GREY)
    label(sl, line, 7.6, 1.15 + i * 0.38, 5.2, 0.35,
          size=Pt(10), color=col, bold=line.startswith("$"))

label(sl, "Anirudh Yadav  ·  github.com/anirudhyadav/aibrd",
      0.5, 7.05, 12, 0.3, size=Pt(9), color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT4)

label(sl, "The Problem", 0.5, 0.25, 12, 0.55,
      size=Pt(32), bold=True, color=WHITE)
label(sl, "Every team loses the thread between the BRD and the code.",
      0.5, 0.9, 12, 0.4, size=Pt(16), color=LIGHT_GREY, italic=True)

problems = [
    ("📋  Requirements drift",
     "BRDs are exported to PDF, shared over email,\nthen forgotten as code evolves in a different direction."),
    ("🔍  No traceability",
     "By release day, no one can confidently answer:\n\"Did we build what the business actually asked for?\""),
    ("🔄  PO updates lost",
     "Mid-sprint requirement changes arrive via Slack or Jira.\nThey never make it back into the spec."),
    ("❌  Test coverage gaps",
     "QA writes tests against what developers built,\nnot against what the BRD required."),
    ("📊  Release anxiety",
     "Release notes are written from memory.\nCompliance sign-off has no traceable evidence."),
    ("⏳  Stale specs",
     "After 3 months, no one trusts the spec.\nAfter 6 months, it's ignored entirely."),
]

for i, (title, body) in enumerate(problems):
    col = i % 3
    row = i // 3
    left = 0.4 + col * 4.3
    top  = 1.5 + row * 2.4
    card(sl, left, top, 4.1, 2.2, title, body.split("\n"),
         accent_color=[ACCENT4, ACCENT3, ACCENT][i % 3])


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — How It Works (pipeline)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT2)

label(sl, "How It Works", 0.5, 0.25, 12, 0.55,
      size=Pt(32), bold=True, color=WHITE)
label(sl, "One command. Every role served. The spec lives in git.",
      0.5, 0.88, 12, 0.35, size=Pt(14), color=LIGHT_GREY, italic=True)

# Pipeline steps
steps = [
    (ACCENT,  "1", "BRD Input",       "PDF · DOCX · MD\nConfluence page"),
    (ACCENT2, "2", "Parse & Chunk",    "Token-aware splitter\nPreserves context"),
    (ACCENT3, "3", "LLM Extract",      "Actors · Flows · Rules\nAcceptance Criteria"),
    (ACCENT,  "4", "Assign IDs",       "PAY-BF-001 stable\nNever reused"),
    (ACCENT2, "5", "Generate Outputs", "CONTEXT.md · RTM\nTests · Reports"),
    (ACCENT3, "6", "Commit to Git",    "Version-controlled\nLiving spec"),
]

for i, (col, num, title, body) in enumerate(steps):
    left = 0.4 + i * 2.15
    # circle
    circ = sl.shapes.add_shape(9, Inches(left + 0.7), Inches(1.35), Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = col
    circ.line.fill.background()
    label(sl, num, left + 0.72, 1.38, 0.65, 0.5,
          size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # arrow
    if i < len(steps) - 1:
        label(sl, "→", left + 1.5, 1.52, 0.6, 0.35,
              size=Pt(20), color=MID_GREY, align=PP_ALIGN.CENTER)
    box(sl, left, 2.2, 2.0, 1.6, fill_color=CARD_BG)
    label(sl, title, left + 0.1, 2.28, 1.85, 0.35,
          size=Pt(12), bold=True, color=WHITE)
    label(sl, body, left + 0.1, 2.65, 1.85, 0.9,
          size=Pt(10), color=LIGHT_GREY)

# Downstream usage row
label(sl, "What each role does with the living spec:", 0.5, 4.05, 12, 0.35,
      size=Pt(13), color=LIGHT_GREY, italic=True)

roles = [
    (ACCENT,  "Developer",      "@aibrd tasks\n@aibrd coverage"),
    (ACCENT2, "QA / Tester",    "aibrd: Generate\nTest Cases"),
    (ACCENT3, "Scrum Master",   "aibrd: Generate\nSprint Feed"),
    (ACCENT4, "Architect",      "aibrd: Derive\nAPI Contracts"),
    (ACCENT,  "Release Mgr",    "aibrd: Release\nNotes"),
    (ACCENT2, "Product Owner",  "aibrd: PO\nProgress Report"),
    (ACCENT3, "Compliance",     "aibrd: Map\nFrameworks"),
]

for i, (col, role, action) in enumerate(roles):
    left = 0.35 + i * 1.85
    box(sl, left, 4.45, 1.75, 1.55, fill_color=CARD_BG)
    box(sl, left, 4.45, 1.75, 0.06, fill_color=col)
    label(sl, role, left + 0.08, 4.55, 1.6, 0.32,
          size=Pt(10), bold=True, color=WHITE)
    label(sl, action, left + 0.08, 4.9, 1.6, 0.8,
          size=Pt(9), color=LIGHT_GREY)

label(sl, "Auto-detected: flat mode (small projects) · modular mode (large, multi-domain projects)",
      0.5, 6.2, 12.3, 0.35, size=Pt(10), color=MID_GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — VS Code Extension
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT)

label(sl, "VS Code Extension", 0.5, 0.25, 10, 0.55,
      size=Pt(32), bold=True, color=WHITE)
label(sl, "GitHub Enterprise + Copilot  ·  No API key  ·  Org-wide VSIX deployment",
      0.5, 0.88, 12, 0.35, size=Pt(13), color=ACCENT, bold=True)

groups = [
    ("Core", ACCENT, [
        ("aibrd: Initialize from BRD",       "Parse PDF/DOCX/MD → full .aibrd/ structure"),
        ("aibrd: Update with new requirement","PO requirement → new IDs → CONTEXT.md"),
        ("aibrd: Generate Test Cases",        "Given/When/Then per AC + boundary per BR"),
        ("aibrd: Generate Release Notes",     "git diff → requirement-mapped release notes"),
        ("aibrd: Show Traceability Matrix",   "RTM tree view in Explorer sidebar"),
        ("aibrd: Show Gap Report",            "Open file vs requirements → coverage gaps"),
    ]),
    ("Quality & Analysis", ACCENT3, [
        ("aibrd: Analyse Change Impact",      "Diff two BRD versions → flag scope drift"),
        ("aibrd: Validate CONTEXT.md",        "Checks cross-refs, dupes, changelogs"),
        ("aibrd: Draft Pull Request Desc.",   "git diff + IDs → traceable PR description"),
    ]),
    ("Delivery Tools", ACCENT2, [
        ("aibrd: Generate Sprint Feed",       "TASK-001 with story points + AC checklist"),
        ("aibrd: Derive API Contracts",       "Flows → OpenAPI 3.0 YAML or Markdown"),
        ("aibrd: Generate PO Report",         "Plain-English progress report for PO"),
        ("aibrd: Map Compliance Frameworks",  "GDPR · WCAG · HIPAA · SOX · PCI-DSS · ISO27001"),
    ]),
    ("Ingestion & Traceability", ACCENT4, [
        ("aibrd: Ingest from Confluence",     "Fetch page + children via REST API"),
        ("aibrd: Check Requirement Staleness","BF-XXX vs git log → 14/30-day drift flags"),
        ("aibrd: Link Requirements to Tests", "Scan test files for ID mentions → coverage %"),
    ]),
]

left_starts = [0.4, 0.4, 7.0, 7.0]
top_starts  = [1.35, 4.1, 1.35, 3.7]
widths      = [6.3, 6.3, 6.0, 6.0]

for gi, (grp_title, col, cmds) in enumerate(groups):
    lft = left_starts[gi]
    top = top_starts[gi]
    w   = widths[gi]
    # header
    box(sl, lft, top, w, 0.32, fill_color=col)
    label(sl, grp_title, lft + 0.1, top + 0.04, w - 0.15, 0.25,
          size=Pt(10), bold=True, color=WHITE)
    for ci, (cmd, desc) in enumerate(cmds):
        row_top = top + 0.35 + ci * 0.42
        box(sl, lft, row_top, w, 0.38, fill_color=CARD_BG)
        label(sl, cmd,  lft + 0.12, row_top + 0.04, w * 0.45, 0.28,
              size=Pt(9), bold=True, color=WHITE)
        label(sl, desc, lft + w * 0.46, row_top + 0.04, w * 0.52, 0.28,
              size=Pt(8.5), color=LIGHT_GREY)

# @aibrd chat note
label(sl, "💬  @aibrd  ·  answer questions · tasks · coverage · rtm  — token-aware, no full-context dump",
      0.4, 7.1, 12.5, 0.28, size=Pt(9.5), color=MID_GREY, bold=False)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Python Library
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT2)

label(sl, "Python Library & CLI", 0.5, 0.25, 10, 0.55,
      size=Pt(32), bold=True, color=WHITE)
label(sl, "Personal repos · CI pipelines · bring-your-own LLM key · import aibrd",
      0.5, 0.88, 12, 0.35, size=Pt(13), color=ACCENT2, bold=True)

# Provider chain
label(sl, "LLM auto-detection:", 0.5, 1.35, 4, 0.3, size=Pt(11), color=LIGHT_GREY)
for i, (prov, key, col) in enumerate([
    ("Anthropic Claude", "ANTHROPIC_API_KEY", ACCENT),
    ("GitHub Models", "GITHUB_TOKEN", ACCENT2),
    ("OpenAI", "OPENAI_API_KEY", ACCENT3),
]):
    box(sl, 0.5 + i * 4.2, 1.68, 4.0, 0.52, fill_color=CARD_BG)
    box(sl, 0.5 + i * 4.2, 1.68, 0.06, 0.52, fill_color=col)
    label(sl, f"{'①②③'[i]}  {prov}", 0.72 + i * 4.2, 1.72, 3.7, 0.25,
          size=Pt(10), bold=True, color=WHITE)
    label(sl, key, 0.72 + i * 4.2, 1.96, 3.7, 0.2,
          size=Pt(8.5), color=ACCENT if i == 0 else (ACCENT2 if i == 1 else ACCENT3))

# CLI columns
cli_groups = [
    ("Core Commands", ACCENT, [
        "aibrd init ./brd.pdf",
        "aibrd update \"PO text\"",
        "aibrd tests [--module slug]",
        "aibrd gaps src/file.py",
        "aibrd release v2.3.0",
    ]),
    ("Quality & Analysis", ACCENT3, [
        "aibrd validate",
        "aibrd pr-draft [--base main]",
        "aibrd change-impact ./brd-v2.pdf",
    ]),
    ("Delivery Tools", ACCENT2, [
        "aibrd sprint [--github-issues]",
        "aibrd api-contracts [--format openapi]",
        "aibrd po-report v2.3.0",
        "aibrd compliance --fw GDPR --fw HIPAA",
    ]),
    ("Ingestion & Traceability", ACCENT4, [
        "aibrd confluence --url … --space ENG",
        "aibrd stale [--module slug]",
        "aibrd test-linkage",
    ]),
]

col_x = [0.4, 0.4, 6.9, 6.9]
col_y = [2.4, 4.55, 2.4, 4.55]
col_w = [6.2, 6.2, 6.1, 6.1]

for gi, (grp, col, cmds) in enumerate(cli_groups):
    lx = col_x[gi]; ty = col_y[gi]; cw = col_w[gi]
    box(sl, lx, ty, cw, 0.3, fill_color=col)
    label(sl, grp, lx + 0.1, ty + 0.04, cw - 0.15, 0.22,
          size=Pt(9.5), bold=True, color=WHITE)
    for ci, cmd in enumerate(cmds):
        ry = ty + 0.33 + ci * 0.37
        box(sl, lx, ry, cw, 0.34, fill_color=CARD_BG)
        label(sl, cmd, lx + 0.12, ry + 0.05, cw - 0.2, 0.25,
              size=Pt(9), color=ACCENT2, bold=True)

label(sl, "import aibrd  ·  All generators, analyzers, parsers available as Python modules",
      0.4, 7.1, 12.5, 0.28, size=Pt(9.5), color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The .aibrd/ Folder + Stable IDs
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT3)

label(sl, "The .aibrd/ Living Specification Folder", 0.5, 0.25, 12, 0.55,
      size=Pt(28), bold=True, color=WHITE)
label(sl, "Commit it to git  ·  Both tools write the same format  ·  Never loses IDs",
      0.5, 0.88, 12, 0.35, size=Pt(13), color=ACCENT3, bold=True)

# Folder tree — flat
flat_lines = [
    ".aibrd/  (flat mode)",
    "├── registry.json       ← ID counter, never reused",
    "├── CONTEXT.md          ← actors · flows · rules · AC",
    "├── index.md            ← traceability matrix",
    "├── ambiguity-report.md",
    "├── conflict-report.md",
    "├── change-impact-report.md",
    "├── compliance-map.md",
    "├── sprint-feed.md",
    "├── staleness-report.md",
    "├── test-linkage-report.md",
    "├── openapi.yaml / api-contracts.md",
    "├── tests/test-cases.md",
    "└── releases/",
    "    ├── v1.0.md",
    "    └── po-report-v1.0.md",
]
box(sl, 0.4, 1.3, 5.8, 5.7, fill_color=RGBColor(0x0F, 0x17, 0x2A))
for i, line in enumerate(flat_lines):
    col = ACCENT3 if i == 0 else (ACCENT if "registry" in line else
          (ACCENT2 if "CONTEXT" in line or "index" in line else LIGHT_GREY))
    label(sl, line, 0.55, 1.4 + i * 0.33, 5.5, 0.3,
          size=Pt(8.5), color=col, bold=(i == 0))

# Modular note
box(sl, 6.4, 1.3, 6.5, 3.1, fill_color=RGBColor(0x0F, 0x17, 0x2A))
modular_lines = [
    ".aibrd/  (modular mode)",
    "├── modules/",
    "│   ├── payments/",
    "│   │   ├── CONTEXT.md   # PAY-BF-001 …",
    "│   │   ├── PAY-openapi.yaml",
    "│   │   └── tests/test-cases.md",
    "│   ├── auth/",
    "│   │   └── CONTEXT.md   # AUTH-BF-001 …",
    "│   └── notifications/",
    "└── shared/",
    "    ├── actors.md        # ACT-001 …",
    "    └── global-rules.md  # GBR-001 …",
]
for i, line in enumerate(modular_lines):
    col = ACCENT3 if i == 0 else LIGHT_GREY
    label(sl, line, 6.55, 1.4 + i * 0.23, 6.2, 0.22,
          size=Pt(8.2), color=col, bold=(i == 0))

# Stable IDs
label(sl, "Stable ID System", 6.4, 4.55, 6.5, 0.35,
      size=Pt(14), bold=True, color=WHITE)
id_rows = [
    ("Type", "Format", "Example", "Generated By"),
    ("Business Flow", "PAY-BF-NNN", "PAY-BF-012", "init / update"),
    ("Business Rule", "PAY-BR-NNN", "PAY-BR-007", "init / update"),
    ("Acceptance Criteria", "PAY-AC-NNN", "PAY-AC-009", "init / update"),
    ("Test Case", "PAY-TC-NNN", "PAY-TC-003", "Generate Tests"),
    ("Release Note", "RN-NNN", "RN-004", "Release Notes"),
    ("Actor", "ACT-NNN", "ACT-002", "init"),
]
row_colors = [MID_GREY] + [CARD_BG] * (len(id_rows) - 1)
for i, (t, fmt, ex, gen) in enumerate(id_rows):
    ry = 4.95 + i * 0.32
    box(sl, 6.4, ry, 6.5, 0.3, fill_color=row_colors[i])
    is_hdr = (i == 0)
    for j, (val, xoff, wid) in enumerate([(t, 0.05, 2.0), (fmt, 2.1, 1.6), (ex, 3.75, 1.3), (gen, 5.1, 1.3)]):
        label(sl, val, 6.4 + xoff, ry + 0.04, wid, 0.24,
              size=Pt(8.5), color=WHITE if is_hdr else LIGHT_GREY, bold=is_hdr)

label(sl, "⚠  Counters only go up — IDs are never reused even if a requirement is deleted",
      0.4, 7.1, 12.5, 0.28, size=Pt(9.5), color=ACCENT3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT)

label(sl, "Architecture", 0.5, 0.25, 12, 0.55,
      size=Pt(32), bold=True, color=WHITE)
label(sl, "All local  ·  No backend  ·  No data leaves the developer machine",
      0.5, 0.88, 12, 0.35, size=Pt(13), color=ACCENT, bold=True)

# Layers
layers = [
    (ACCENT,  "Input Layer",
     ["PDF Parser (pdf-parse)", "DOCX Parser (mammoth)", "Markdown Parser", "Confluence REST Ingester"]),
    (ACCENT2, "LLM Layer",
     ["vscode.lm  (Copilot — no key)", "Anthropic / GitHub Models / OpenAI  (Python)",
      "Token-aware chunker", "callLLM() · callLLMJson<T>()"]),
    (ACCENT3, "Core Layer",
     ["Module Detector", "Extractors: Actors · Flows · Rules · AC",
      "Registry (append-only ID counter)", "Resolver (cross-ref linker)"]),
    (ACCENT4, "Generator Layer",
     ["CONTEXT.md · RTM · Test Cases · Sprint Feed",
      "API Contracts · PO Report · Compliance Map",
      "Release Notes · Ambiguity · Conflict Reports"]),
    (ACCENT,  "Analyzer Layer",
     ["Gap Detector", "Change Impact", "Validator",
      "Stale Detector", "Test Linkage"]),
    (ACCENT2, "Output Layer",
     ["workspace/writer.ts (VS Code)", "workspace/writer.py (Python)",
      ".aibrd/ folder — git-committed"]),
]

for i, (col, title, items) in enumerate(layers):
    lft = 0.4 + (i % 3) * 4.3
    top = 1.35 + (i // 3) * 2.75
    box(sl, lft, top, 4.1, 2.55, fill_color=CARD_BG)
    box(sl, lft, top, 4.1, 0.32, fill_color=col)
    label(sl, title, lft + 0.1, top + 0.05, 3.9, 0.24,
          size=Pt(10), bold=True, color=WHITE)
    for j, item in enumerate(items):
        label(sl, f"  • {item}", lft + 0.08, top + 0.4 + j * 0.48, 3.9, 0.42,
              size=Pt(9), color=LIGHT_GREY)

label(sl, "GitHub Actions reusable workflow  ·  GITHUB_TOKEN + GitHub Models  ·  No separate API key",
      0.4, 7.1, 12.5, 0.28, size=Pt(9.5), color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Deployment & Onboarding
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
accent_bar(sl, ACCENT2)

label(sl, "Deployment & Onboarding", 0.5, 0.25, 12, 0.55,
      size=Pt(32), bold=True, color=WHITE)

# VS Code track
box(sl, 0.4, 0.95, 6.1, 0.38, fill_color=ACCENT)
label(sl, "VS Code Extension — Org Track", 0.55, 0.98, 5.9, 0.3,
      size=Pt(12), bold=True, color=WHITE)

vsc_steps = [
    ("1", "Package", "npm install && npx vsce package\nDistribute .vsix via MDM or VS Code Server"),
    ("2", "Team Opens Repo", "File → Open Folder in VS Code\nCopilot license covers everything"),
    ("3", "Tech Lead Inits", "aibrd: Initialize from BRD\nSelect PDF/DOCX/MD or ingest from Confluence"),
    ("4", "Commit .aibrd/", "git add .aibrd/ && git commit\nSpec is now version-controlled"),
    ("5", "Add CI Workflow", "uses: org/aibrd/.github/workflows/\naibrd-reusable.yml@main"),
    ("6", "Share Playbook", "PLAYBOOK.md covers every role:\nDev · QA · Lead · Release Mgr · PO"),
]
for i, (num, title, body) in enumerate(vsc_steps):
    lft = 0.4 + (i % 3) * 2.05
    top = 1.45 + (i // 3) * 1.75
    box(sl, lft, top, 1.95, 1.65, fill_color=CARD_BG)
    circ = sl.shapes.add_shape(9, Inches(lft + 0.7), Inches(top + 0.08),
                                Inches(0.52), Inches(0.52))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    label(sl, num, lft + 0.72, top + 0.1, 0.5, 0.35,
          size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(sl, title, lft + 0.08, top + 0.66, 1.8, 0.3,
          size=Pt(9.5), bold=True, color=WHITE)
    label(sl, body, lft + 0.08, top + 0.96, 1.8, 0.65,
          size=Pt(8), color=LIGHT_GREY)

# Python track
box(sl, 6.8, 0.95, 6.1, 0.38, fill_color=ACCENT2)
label(sl, "Python Library — Personal / CI Track", 6.95, 0.98, 5.9, 0.3,
      size=Pt(12), bold=True, color=WHITE)

py_steps = [
    ("1", "Install", "pip install -e pythonlibrary/\naibrd --help"),
    ("2", "Set LLM Key", "export ANTHROPIC_API_KEY=...\nor GITHUB_TOKEN (free)"),
    ("3", "Init Project", "aibrd init ./brd.pdf\n.aibrd/ generated"),
    ("4", "Commit .aibrd/", "git add .aibrd/ && git commit\nSpec in git"),
    ("5", "Daily Usage", "aibrd sprint  ·  aibrd gaps\naibrd validate  ·  aibrd stale"),
    ("6", "CI Integration", "aibrd validate → exit 1 on fail\naibrd stale → exit 1 if stale"),
]
for i, (num, title, body) in enumerate(py_steps):
    lft = 6.8 + (i % 3) * 2.05
    top = 1.45 + (i // 3) * 1.75
    box(sl, lft, top, 1.95, 1.65, fill_color=CARD_BG)
    circ = sl.shapes.add_shape(9, Inches(lft + 0.7), Inches(top + 0.08),
                                Inches(0.52), Inches(0.52))
    circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT2
    circ.line.fill.background()
    label(sl, num, lft + 0.72, top + 0.1, 0.5, 0.35,
          size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    label(sl, title, lft + 0.08, top + 0.66, 1.8, 0.3,
          size=Pt(9.5), bold=True, color=WHITE)
    label(sl, body, lft + 0.08, top + 0.96, 1.8, 0.65,
          size=Pt(8), color=LIGHT_GREY)

label(sl, "Scale: same .aibrd/ format across 100s of Python and Java repos — one reusable workflow handles all",
      0.4, 7.1, 12.5, 0.28, size=Pt(9.5), color=MID_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — COMPLETE FEATURE LIST (last slide)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
bg(sl)
box(sl, 0, 0, 13.33, 7.5, fill_color=RGBColor(0x09, 0x0E, 0x18))  # extra dark

# Header strip
box(sl, 0, 0, 13.33, 0.85, fill_color=DARK_BG)
label(sl, "Complete Feature List", 0.4, 0.08, 9, 0.55,
      size=Pt(26), bold=True, color=WHITE)
label(sl, "aibrd v0.2.0", 11.3, 0.12, 2, 0.35,
      size=Pt(13), color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
label(sl, "16 VS Code commands  ·  15 Python CLI commands  ·  Modular + Flat mode",
      0.4, 0.55, 12.5, 0.28, size=Pt(9), color=MID_GREY)

# 4 columns x 3 batches + meta
columns = [
    {
        "title": "Core",
        "color": ACCENT,
        "items": [
            ("aibrd: Initialize from BRD",     "Parse PDF/DOCX/MD → actors, flows, rules, AC, IDs"),
            ("aibrd: Update",                   "PO requirement → stable IDs → CONTEXT.md append"),
            ("aibrd: Generate Test Cases",      "Gherkin per AC + boundary per BR"),
            ("aibrd: Generate Release Notes",   "git diff → requirement-ID mapped release notes"),
            ("aibrd: Show Traceability Matrix", "RTM tree view in VS Code Explorer sidebar"),
            ("aibrd: Show Gap Report",          "Open file vs CONTEXT.md → coverage gaps"),
            ("@aibrd (Copilot Chat)",           "tasks · coverage · analyze ID · rtm queries"),
        ],
    },
    {
        "title": "Quality & Analysis",
        "color": ACCENT3,
        "items": [
            ("aibrd: Analyse Change Impact",    "Diff two BRD versions → scope drift, new, removed"),
            ("aibrd: Validate CONTEXT.md",      "Cross-refs, duplicate IDs, changelogs — pure logic"),
            ("aibrd: Draft Pull Request Desc.", "git diff + requirements → traceable PR description"),
            ("aibrd validate  (CLI)",           "Exit 1 on failure — safe for CI"),
            ("aibrd pr-draft  (CLI)",           "Base branch configurable, prints to stdout"),
            ("aibrd change-impact  (CLI)",      "File input, saves change-impact-report.md"),
        ],
    },
    {
        "title": "Delivery Tools",
        "color": ACCENT2,
        "items": [
            ("aibrd: Generate Sprint Feed",     "TASK-001 · story points · priority · AC checklist"),
            ("aibrd: Derive API Contracts",     "Business flows → OpenAPI 3.0 YAML or Markdown"),
            ("aibrd: Generate PO Report",       "Plain-English, no IDs visible — for PO sign-off"),
            ("aibrd: Map Compliance",           "GDPR · WCAG · HIPAA · SOX · PCI-DSS · ISO27001"),
            ("aibrd sprint  (CLI)",             "--github-issues for GitHub Issues API JSON"),
            ("aibrd api-contracts  (CLI)",      "--format openapi or markdown, per module"),
            ("aibrd po-report  (CLI)",          "git range configurable, saves to releases/"),
            ("aibrd compliance  (CLI)",         "--fw flag repeatable for multiple frameworks"),
        ],
    },
    {
        "title": "Ingestion & Traceability",
        "color": ACCENT4,
        "items": [
            ("aibrd: Ingest from Confluence",   "Cloud (email+token) or Server/DC (PAT) via REST"),
            ("aibrd: Check Requirement Staleness", "BF-XXX vs git log → 🔴 Stale 🟡 Drifting 🟢 Ok"),
            ("aibrd: Link Requirements to Tests",  "Scan test files for ID mentions → coverage %"),
            ("aibrd confluence  (CLI)",         "CONFLUENCE_TOKEN env var supported"),
            ("aibrd stale  (CLI)",              "Exit 1 if stale — CI-safe"),
            ("aibrd test-linkage  (CLI)",       "Per-module or all, saves test-linkage-report.md"),
        ],
    },
]

col_x = [0.3, 3.65, 7.0, 10.05]
col_w = [3.25, 3.25, 2.98, 3.2]

for ci, col_data in enumerate(columns):
    lx = col_x[ci]
    cw = col_w[ci]
    col = col_data["color"]

    # Column header
    box(sl, lx, 0.88, cw, 0.35, fill_color=col)
    label(sl, col_data["title"], lx + 0.08, 0.9, cw - 0.1, 0.28,
          size=Pt(10), bold=True, color=WHITE)

    for ri, (feat, desc) in enumerate(col_data["items"]):
        ry = 1.26 + ri * 0.79
        is_cli = "(CLI)" in feat
        row_col = RGBColor(0x14, 0x1E, 0x30) if is_cli else CARD_BG
        box(sl, lx, ry, cw, 0.75, fill_color=row_col)
        # small left accent
        box(sl, lx, ry, 0.04, 0.75, fill_color=col)
        # CLI badge
        if is_cli:
            box(sl, lx + cw - 0.62, ry + 0.04, 0.58, 0.22, fill_color=ACCENT2)
            label(sl, "CLI", lx + cw - 0.6, ry + 0.06, 0.56, 0.18,
                  size=Pt(7), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        label(sl, feat.replace(" (CLI)", ""),
              lx + 0.1, ry + 0.04, cw - 0.15, 0.28,
              size=Pt(8.5), bold=True, color=WHITE)
        label(sl, desc, lx + 0.1, ry + 0.34, cw - 0.15, 0.38,
              size=Pt(7.5), color=LIGHT_GREY)

# Footer
box(sl, 0, 7.2, 13.33, 0.3, fill_color=DARK_BG)
label(sl, "Both tools write the same .aibrd/ format  ·  Flat + Modular mode auto-detected  ·"
         "  IDs never reused  ·  No API key for VS Code  ·  Commit .aibrd/ to git",
      0.3, 7.22, 12.8, 0.25, size=Pt(7.8), color=MID_GREY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/anirudhyadav/Documents/GenAI_Architect/Claude/aibrd/aibrd-product-deck.pptx"
prs.save(out)
print(f"Saved → {out}")
print(f"Slides: {len(prs.slides)}")
